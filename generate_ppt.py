from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Color Palette ──
PRIMARY = RGBColor(37, 99, 235)
DARK = RGBColor(15, 23, 42)
WHITE = RGBColor(255, 255, 255)
MUTED = RGBColor(100, 116, 139)
SUCCESS = RGBColor(16, 185, 129)
DANGER = RGBColor(239, 68, 68)
LIGHT_BG = RGBColor(241, 245, 249)
ACCENT = RGBColor(99, 102, 241)
CARD_BG = RGBColor(248, 250, 252)

def add_bg(slide, color=DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_gradient_bg(slide):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = RGBColor(15, 23, 42)
    fill.gradient_stops[0].position = 0
    fill.gradient_stops[1].color.rgb = RGBColor(30, 41, 82)
    fill.gradient_stops[1].position = 1.0

def add_shape(slide, left, top, width, height, color, radius=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    if radius:
        shape.adjustments[0] = radius
    return shape

def add_textbox(slide, left, top, width, height, text="", font_size=18, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_bullet_list(slide, left, top, width, height, items, font_size=16, color=WHITE, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.space_after = spacing
        p.level = 0
    return txBox

def add_icon_card(slide, left, top, width, height, icon_text, title, desc, accent_color=PRIMARY):
    card = add_shape(slide, left, top, width, height, RGBColor(30, 41, 59), 0.05)
    # Icon circle
    icon_shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left + Inches(0.4), top + Inches(0.4), Inches(0.7), Inches(0.7))
    icon_shape.fill.solid()
    icon_shape.fill.fore_color.rgb = accent_color
    icon_shape.line.fill.background()
    icon_tf = icon_shape.text_frame
    icon_tf.paragraphs[0].text = icon_text
    icon_tf.paragraphs[0].font.size = Pt(22)
    icon_tf.paragraphs[0].font.color.rgb = WHITE
    icon_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    icon_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    # Title
    add_textbox(slide, left + Inches(0.3), top + Inches(1.3), width - Inches(0.6), Inches(0.5), title, 17, WHITE, True)
    # Description
    add_textbox(slide, left + Inches(0.3), top + Inches(1.8), width - Inches(0.6), Inches(1.2), desc, 13, MUTED)


# ════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank
add_gradient_bg(slide)

# Decorative accent bar
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), PRIMARY)

# Badge
badge = add_shape(slide, Inches(1.5), Inches(1.5), Inches(3), Inches(0.5), PRIMARY, 0.1)
badge_tf = badge.text_frame
badge_tf.paragraphs[0].text = "FINAL YEAR PROJECT"
badge_tf.paragraphs[0].font.size = Pt(13)
badge_tf.paragraphs[0].font.color.rgb = WHITE
badge_tf.paragraphs[0].font.bold = True
badge_tf.paragraphs[0].font.name = "Calibri"
badge_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
badge_tf.vertical_anchor = MSO_ANCHOR.MIDDLE

# Title
add_textbox(slide, Inches(1.5), Inches(2.3), Inches(8), Inches(1.2), "Smart ERP", 54, WHITE, True)
add_textbox(slide, Inches(1.5), Inches(3.3), Inches(8), Inches(0.8), "Biometric Attendance Management System", 28, PRIMARY, False)
add_textbox(slide, Inches(1.5), Inches(4.2), Inches(8), Inches(0.6), "AI-Powered Face Recognition  •  QR Authentication  •  Real-Time Analytics", 16, MUTED)

# Divider line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.5), Inches(5.2), Inches(3), Pt(2))
line.fill.solid()
line.fill.fore_color.rgb = PRIMARY
line.line.fill.background()

# Team info
add_textbox(slide, Inches(1.5), Inches(5.5), Inches(5), Inches(0.4), "Presented by: Prajith Kumar Reddy", 16, WHITE, True)
add_textbox(slide, Inches(1.5), Inches(5.9), Inches(5), Inches(0.4), "KL University  •  2026", 14, MUTED)

# Right side decorative elements
circle1 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(10), Inches(1), Inches(3), Inches(3))
circle1.fill.solid()
circle1.fill.fore_color.rgb = RGBColor(37, 99, 235)
circle1.fill.fore_color.brightness = 0.8
circle1.line.fill.background()

circle2 = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(9), Inches(4), Inches(2), Inches(2))
circle2.fill.solid()
circle2.fill.fore_color.rgb = RGBColor(99, 102, 241)
circle2.fill.fore_color.brightness = 0.85
circle2.line.fill.background()


# ════════════════════════════════════════════════════════════
# SLIDE 2: PROBLEM STATEMENT
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), DANGER)

add_textbox(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.5), "PROBLEM STATEMENT", 14, DANGER, True)
add_textbox(slide, Inches(1), Inches(1.0), Inches(10), Inches(0.8), "Why Traditional Attendance Fails", 40, WHITE, True)

problems = [
    ("⏱️", "Time Consuming", "Manual roll calls waste 10-15 minutes of each lecture, reducing actual teaching time significantly.", DANGER),
    ("🔄", "Proxy Attendance", "Students mark attendance for absent friends using buddy systems, manipulating paper-based records.", RGBColor(245, 158, 11)),
    ("📊", "No Real-Time Data", "Faculty cannot track trends, patterns, or at-risk students until end of semester report generation.", ACCENT),
    ("📝", "Error Prone Records", "Paper registers are prone to data entry errors, damage, loss, and difficult to digitize for analytics.", MUTED),
]

for i, (icon, title, desc, color) in enumerate(problems):
    x = Inches(0.8 + (i % 2) * 6.2)
    y = Inches(2.2 + (i // 2) * 2.6)
    add_icon_card(slide, x, y, Inches(5.6), Inches(2.3), icon, title, desc, color)


# ════════════════════════════════════════════════════════════
# SLIDE 3: SOLUTION OVERVIEW
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), SUCCESS)

add_textbox(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.5), "OUR SOLUTION", 14, SUCCESS, True)
add_textbox(slide, Inches(1), Inches(1.0), Inches(10), Inches(0.8), "Smart ERP: Intelligent Attendance", 40, WHITE, True)

steps = [
    ("1", "Employee\nGenerates QR", "A time-limited\ncryptographic QR token\nis generated for\neach session"),
    ("2", "Student\nScans QR", "Student scans the\nQR code from their\nmobile browser to\nopen the portal"),
    ("3", "Face\nVerification", "AI-powered face\nrecognition verifies\nstudent identity\nagainst registered data"),
    ("4", "Attendance\nMarked", "Records are saved\nwith timestamp and\nupdated in real-time\nacross all dashboards"),
]

for i, (num, title, desc) in enumerate(steps):
    x = Inches(0.6 + i * 3.2)
    y = Inches(2.5)
    
    # Number circle
    circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.85), y, Inches(0.9), Inches(0.9))
    circle.fill.solid()
    circle.fill.fore_color.rgb = PRIMARY
    circle.line.fill.background()
    tf = circle.text_frame
    tf.paragraphs[0].text = num
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    # Arrow connector (except last)  
    if i < 3:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.6), y + Inches(0.25), Inches(0.6), Inches(0.4))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = MUTED
        arrow.line.fill.background()
    
    # Card below
    card = add_shape(slide, x, y + Inches(1.2), Inches(2.8), Inches(3.0), RGBColor(30, 41, 59), 0.05)
    add_textbox(slide, x + Inches(0.2), y + Inches(1.4), Inches(2.4), Inches(0.7), title, 18, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.2), y + Inches(2.1), Inches(2.4), Inches(1.8), desc, 13, MUTED, False, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 4: TECH STACK
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT)

add_textbox(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.5), "TECHNOLOGY STACK", 14, ACCENT, True)
add_textbox(slide, Inches(1), Inches(1.0), Inches(10), Inches(0.8), "Modern, Scalable Architecture", 40, WHITE, True)

stacks = [
    ("Frontend", [
        "React 19 + Vite 8",
        "React Router v7",
        "face-api.js (TensorFlow.js)",
        "html5-qrcode Scanner",
        "Lucide React Icons",
    ], PRIMARY),
    ("Backend", [
        "Node.js + Express",
        "JWT Authentication",
        "Euclidean Distance Matching",
        "RESTful API Architecture",
        "CORS Security Layer",
    ], SUCCESS),
    ("Database", [
        "MongoDB Atlas (Cloud)",
        "Mongoose ODM",
        "Indexed Queries",
        "128-D Face Descriptors",
        "Attendance Schema",
    ], RGBColor(245, 158, 11)),
    ("Deployment", [
        "Vercel (Serverless)",
        "Separate Client/Server",
        "Auto CI/CD from GitHub",
        "Environment Variables",
        "Edge Network CDN",
    ], DANGER),
]

for i, (title, items, color) in enumerate(stacks):
    x = Inches(0.6 + i * 3.2)
    y = Inches(2.2)
    
    card = add_shape(slide, x, y, Inches(2.9), Inches(4.5), RGBColor(30, 41, 59), 0.04)
    
    # Colored top bar
    bar = add_shape(slide, x, y, Inches(2.9), Inches(0.6), color, 0.04)
    bar_tf = bar.text_frame
    bar_tf.paragraphs[0].text = title
    bar_tf.paragraphs[0].font.size = Pt(18)
    bar_tf.paragraphs[0].font.color.rgb = WHITE
    bar_tf.paragraphs[0].font.bold = True
    bar_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    bar_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    add_bullet_list(slide, x + Inches(0.25), y + Inches(0.9), Inches(2.4), Inches(3.3), 
                    [f"▸  {item}" for item in items], 13, RGBColor(203, 213, 225), Pt(10))


# ════════════════════════════════════════════════════════════
# SLIDE 5: SYSTEM ARCHITECTURE
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), PRIMARY)

add_textbox(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.5), "SYSTEM ARCHITECTURE", 14, PRIMARY, True)
add_textbox(slide, Inches(1), Inches(1.0), Inches(10), Inches(0.8), "Three-Tier Role-Based Architecture", 40, WHITE, True)

roles = [
    ("👨‍💼", "Admin Portal", [
        "View system metrics",
        "Manage students & employees",
        "Configure fee structures",
        "Global enrollment control",
    ], DANGER),
    ("👨‍🏫", "Employee Portal", [
        "Generate QR attendance sessions",
        "Real-time student tracking",
        "Manual override attendance",
        "Session finalization controls",
    ], PRIMARY),
    ("🎓", "Student Portal", [
        "Register biometric face ID",
        "Scan QR → Face verification",
        "View attendance %",
        "Track session history",
    ], SUCCESS),
]

for i, (icon, title, features, color) in enumerate(roles):
    x = Inches(0.8 + i * 4.1)
    y = Inches(2.2)
    w = Inches(3.7)
    
    card = add_shape(slide, x, y, w, Inches(4.8), RGBColor(30, 41, 59), 0.04)

    # Icon
    icon_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(1.2), y + Inches(0.4), Inches(1.2), Inches(1.2))
    icon_circle.fill.solid()
    icon_circle.fill.fore_color.rgb = color
    icon_circle.line.fill.background()
    tf = icon_circle.text_frame
    tf.paragraphs[0].text = icon
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE

    add_textbox(slide, x + Inches(0.2), y + Inches(1.8), w - Inches(0.4), Inches(0.5), title, 20, WHITE, True, PP_ALIGN.CENTER)
    
    # Divider
    div = add_shape(slide, x + Inches(0.8), y + Inches(2.4), Inches(2), Pt(2), color)
    
    add_bullet_list(slide, x + Inches(0.3), y + Inches(2.7), w - Inches(0.6), Inches(2.0),
                    [f"✦  {f}" for f in features], 13, RGBColor(203, 213, 225), Pt(8))


# ════════════════════════════════════════════════════════════
# SLIDE 6: SECURITY FEATURES
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), DANGER)

add_textbox(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.5), "SECURITY LAYER", 14, DANGER, True)
add_textbox(slide, Inches(1), Inches(1.0), Inches(10), Inches(0.8), "Multi-Layer Anti-Fraud System", 40, WHITE, True)

security = [
    ("🔒", "Face Duplication Block", "Same face cannot be registered under multiple student accounts. System cross-checks against all existing descriptors.", SUCCESS),
    ("🛡️", "Proxy Attendance Detection", "During marking, the face is verified as a closer match to the logged-in student than to ANY other registered student.", PRIMARY),
    ("⏰", "Time-Limited QR Tokens", "QR session tokens expire in 2 minutes using JWT cryptographic signing. Cannot be reused or forged.", RGBColor(245, 158, 11)),
    ("🔐", "JWT Authentication", "All API endpoints are protected with Bearer token authentication. Tokens are role-specific (admin/employee/student).", ACCENT),
    ("📏", "Euclidean Distance", "128-dimensional face descriptors compared using mathematical distance. Threshold: 0.55 for high precision.", DANGER),
    ("🌐", "CORS Protection", "Backend only accepts requests from the configured frontend domain. Cross-origin attacks are blocked.", MUTED),
]

for i, (icon, title, desc, color) in enumerate(security):
    x = Inches(0.6 + (i % 3) * 4.2)
    y = Inches(2.2 + (i // 3) * 2.6)
    add_icon_card(slide, x, y, Inches(3.8), Inches(2.3), icon, title, desc, color)


# ════════════════════════════════════════════════════════════
# SLIDE 7: FACE RECOGNITION WORKFLOW
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT)

add_textbox(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.5), "FACE RECOGNITION", 14, ACCENT, True)
add_textbox(slide, Inches(1), Inches(1.0), Inches(10), Inches(0.8), "AI-Powered Biometric Pipeline", 40, WHITE, True)

pipeline = [
    ("Camera\nCapture", "Browser's\ngetUserMedia API\ncaptures live\nvideo feed"),
    ("Face\nDetection", "TinyFaceDetector\nlocates face in\nframe with\nbounding box"),
    ("Landmark\nExtraction", "68-point facial\nlandmark mapping\nfor geometric\nanalysis"),
    ("Descriptor\nGeneration", "128-dimensional\nfloat vector\ngenerated as\nface fingerprint"),
    ("Distance\nMatching", "Euclidean distance\ncompared against\nstored descriptor\n(threshold: 0.55)"),
]

for i, (title, desc) in enumerate(pipeline):
    x = Inches(0.4 + i * 2.55)
    y = Inches(2.5)
    
    # Arrow
    if i < 4:
        arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + Inches(2.1), y + Inches(1.5), Inches(0.5), Inches(0.35))
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = PRIMARY
        arrow.line.fill.background()
    
    card = add_shape(slide, x, y, Inches(2.2), Inches(3.5), RGBColor(30, 41, 59), 0.05)
    
    # Step number
    num_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + Inches(0.75), y + Inches(0.3), Inches(0.65), Inches(0.65))
    num_circle.fill.solid()
    num_circle.fill.fore_color.rgb = ACCENT
    num_circle.line.fill.background()
    tf = num_circle.text_frame
    tf.paragraphs[0].text = str(i + 1)
    tf.paragraphs[0].font.size = Pt(22)
    tf.paragraphs[0].font.color.rgb = WHITE
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    add_textbox(slide, x + Inches(0.15), y + Inches(1.1), Inches(1.9), Inches(0.7), title, 15, WHITE, True, PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.15), y + Inches(1.9), Inches(1.9), Inches(1.4), desc, 12, MUTED, False, PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════
# SLIDE 8: REAL-TIME FEATURES
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), SUCCESS)

add_textbox(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.5), "REAL-TIME FEATURES", 14, SUCCESS, True)
add_textbox(slide, Inches(1), Inches(1.0), Inches(10), Inches(0.8), "Live Dashboard Updates", 40, WHITE, True)

features = [
    ("📊", "Live Attendance Percentage", "Student dashboards show a circular progress indicator with color-coded attendance percentage (green ≥75%, yellow ≥50%, red <50%). Updates every 3 seconds.", SUCCESS),
    ("👥", "Session Roster Matrix", "Employee dashboard shows real-time numbered list of verified students with roll numbers as they scan their faces during active QR sessions.", PRIMARY),
    ("📱", "Mobile QR → Camera Flow", "Students scan QR on phone → Login (if needed) → Camera auto-launches for instant face verification → Attendance marked automatically.", ACCENT),
    ("🔄", "Auto-Refresh Polling", "All dashboards poll the backend every 3 seconds. When a student marks attendance, it reflects instantly across admin, employee, and student views.", RGBColor(245, 158, 11)),
]

for i, (icon, title, desc, color) in enumerate(features):
    x = Inches(0.6 + (i % 2) * 6.2)
    y = Inches(2.2 + (i // 2) * 2.6)
    add_icon_card(slide, x, y, Inches(5.6), Inches(2.3), icon, title, desc, color)


# ════════════════════════════════════════════════════════════
# SLIDE 9: DATABASE DESIGN
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), RGBColor(245, 158, 11))

add_textbox(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.5), "DATABASE DESIGN", 14, RGBColor(245, 158, 11), True)
add_textbox(slide, Inches(1), Inches(1.0), Inches(10), Inches(0.8), "MongoDB Collections Schema", 40, WHITE, True)

schemas = [
    ("Student", [
        "name: String",
        "rollNumber: String (unique)",
        "password: String (hashed)",
        "registeredFace: Boolean",
        "faceDescriptor: [Number] (128-D)",
    ]),
    ("Employee", [
        "name: String",
        "username: String (unique)",
        "password: String (hashed)",
        "department: String",
        "role: 'employee'",
    ]),
    ("Attendance", [
        "student: ObjectId → Student",
        "date: Date (required)",
        "status: 'Present' | 'Absent'",
        "sessionToken: String (QR link)",
        "Index: (student, date) unique",
    ]),
]

for i, (name, fields) in enumerate(schemas):
    x = Inches(0.6 + i * 4.2)
    y = Inches(2.3)
    
    card = add_shape(slide, x, y, Inches(3.8), Inches(4.5), RGBColor(30, 41, 59), 0.04)
    
    header = add_shape(slide, x, y, Inches(3.8), Inches(0.7), RGBColor(245, 158, 11), 0.04)
    htf = header.text_frame
    htf.paragraphs[0].text = f"📦  {name}"
    htf.paragraphs[0].font.size = Pt(20)
    htf.paragraphs[0].font.color.rgb = WHITE
    htf.paragraphs[0].font.bold = True
    htf.paragraphs[0].alignment = PP_ALIGN.CENTER
    htf.vertical_anchor = MSO_ANCHOR.MIDDLE
    
    add_bullet_list(slide, x + Inches(0.25), y + Inches(1.0), Inches(3.3), Inches(3.2),
                    [f"▸  {f}" for f in fields], 13, RGBColor(203, 213, 225), Pt(12))


# ════════════════════════════════════════════════════════════
# SLIDE 10: DEPLOYMENT
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), PRIMARY)

add_textbox(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.5), "DEPLOYMENT & HOSTING", 14, PRIMARY, True)
add_textbox(slide, Inches(1), Inches(1.0), Inches(10), Inches(0.8), "Vercel Cloud Infrastructure", 40, WHITE, True)

# Left panel
left_card = add_shape(slide, Inches(0.8), Inches(2.3), Inches(5.5), Inches(4.5), RGBColor(30, 41, 59), 0.04)
add_textbox(slide, Inches(1.2), Inches(2.5), Inches(4.5), Inches(0.5), "Architecture Overview", 22, WHITE, True)
add_bullet_list(slide, Inches(1.2), Inches(3.2), Inches(4.8), Inches(3.5), [
    "▸  Monorepo: /client + /server directories",
    "▸  Two separate Vercel projects from same repo",
    "▸  Frontend: Static build (Vite → dist/)",
    "▸  Backend: Serverless functions (Node.js)",
    "▸  SPA Routing: vercel.json rewrites",
    "▸  Auto-deploy on every git push",
    "▸  Environment variables for secrets",
], 14, RGBColor(203, 213, 225), Pt(10))

# Right panel - env vars
right_card = add_shape(slide, Inches(6.8), Inches(2.3), Inches(5.5), Inches(4.5), RGBColor(30, 41, 59), 0.04)
add_textbox(slide, Inches(7.2), Inches(2.5), Inches(4.5), Inches(0.5), "Environment Variables", 22, WHITE, True)

env_items = [
    ("Backend", "MONGODB_URI", "Atlas connection string"),
    ("Backend", "JWT_SECRET", "Token signing key"),
    ("Backend", "CORS_ORIGIN", "Frontend URL"),
    ("Frontend", "VITE_API_URL", "Backend API URL"),
]

for i, (where, key, desc) in enumerate(env_items):
    yy = Inches(3.3) + Inches(i * 0.8)
    badge_color = PRIMARY if where == "Backend" else SUCCESS
    b = add_shape(slide, Inches(7.2), yy, Inches(1.1), Inches(0.35), badge_color, 0.15)
    btf = b.text_frame
    btf.paragraphs[0].text = where
    btf.paragraphs[0].font.size = Pt(10)
    btf.paragraphs[0].font.color.rgb = WHITE
    btf.paragraphs[0].font.bold = True
    btf.paragraphs[0].alignment = PP_ALIGN.CENTER
    btf.vertical_anchor = MSO_ANCHOR.MIDDLE
    add_textbox(slide, Inches(8.5), yy - Inches(0.02), Inches(3.5), Inches(0.25), key, 13, WHITE, True)
    add_textbox(slide, Inches(8.5), yy + Inches(0.22), Inches(3.5), Inches(0.25), desc, 11, MUTED)


# ════════════════════════════════════════════════════════════
# SLIDE 11: FUTURE SCOPE
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), ACCENT)

add_textbox(slide, Inches(1), Inches(0.5), Inches(5), Inches(0.5), "FUTURE SCOPE", 14, ACCENT, True)
add_textbox(slide, Inches(1), Inches(1.0), Inches(10), Inches(0.8), "Planned Enhancements", 40, WHITE, True)

future = [
    ("📍", "Geolocation Fencing", "Verify student is physically inside the classroom using GPS coordinates before allowing attendance marking.", SUCCESS),
    ("📧", "Email Notifications", "Automated alerts to parents/students when attendance drops below 75% threshold.", PRIMARY),
    ("📈", "Analytics Dashboard", "Heatmaps, trend graphs, and predictive analytics for identifying at-risk students.", RGBColor(245, 158, 11)),
    ("🔗", "LMS Integration", "Connect with university Learning Management Systems for unified student data.", ACCENT),
    ("🤖", "Liveness Detection", "Anti-spoofing layer to detect if a real person is present vs a photo/video.", DANGER),
    ("📱", "Native Mobile App", "Dedicated iOS/Android app for faster camera access and push notifications.", MUTED),
]

for i, (icon, title, desc, color) in enumerate(future):
    x = Inches(0.6 + (i % 3) * 4.2)
    y = Inches(2.2 + (i // 3) * 2.6)
    add_icon_card(slide, x, y, Inches(3.8), Inches(2.3), icon, title, desc, color)


# ════════════════════════════════════════════════════════════
# SLIDE 12: THANK YOU
# ════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_gradient_bg(slide)
add_shape(slide, Inches(0), Inches(0), Inches(0.15), Inches(7.5), PRIMARY)

# Large centered circle
big_circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(5.1), Inches(1.2), Inches(3), Inches(3))
big_circle.fill.solid()
big_circle.fill.fore_color.rgb = PRIMARY
big_circle.line.fill.background()
tf = big_circle.text_frame
tf.paragraphs[0].text = "🎓"
tf.paragraphs[0].font.size = Pt(60)
tf.paragraphs[0].alignment = PP_ALIGN.CENTER
tf.vertical_anchor = MSO_ANCHOR.MIDDLE

add_textbox(slide, Inches(2), Inches(4.5), Inches(9), Inches(1), "Thank You!", 52, WHITE, True, PP_ALIGN.CENTER)
add_textbox(slide, Inches(2), Inches(5.5), Inches(9), Inches(0.5), "Questions & Live Demo", 24, PRIMARY, False, PP_ALIGN.CENTER)
add_textbox(slide, Inches(2), Inches(6.2), Inches(9), Inches(0.5), "Smart ERP  •  Prajith Kumar Reddy  •  KL University  •  2026", 14, MUTED, False, PP_ALIGN.CENTER)


# ── Save ──
output = r"c:\Users\praji\Downloads\Smart_Attendence\Smart_ERP_Presentation.pptx"
prs.save(output)
print(f"✅ Presentation saved to: {output}")
print(f"📊 Total slides: {len(prs.slides)}")
